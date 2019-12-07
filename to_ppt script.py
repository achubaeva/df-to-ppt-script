from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.util import Pt

'''Create a powerpoint presentation with title page and following pages with tables from dataframe.'''

def create_ppt(title_str, subtitle_str, dictionary_of_dfs, filepath):
    '''Title + subtitle page.'''
    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title

    title.text = title_str


    
    
    
    #subtitle = slide.placeholders[1]
    
    
    #subtitle.text = subtitle_str
    
    '''For the subsequent pages'''
    # Round numbers if necessary
    #df = df.round(2)
    
    for i in dictionary_of_dfs:
        df = dictionary_of_dfs[i]
        df.round(1)
        # First establish title only slide
        title_only_slide_layout = prs.slide_layouts[5]
        # Add slide to presentation
        slide = prs.slides.add_slide(title_only_slide_layout)
        #title = slide.placeholders[0] #new
        shapes = slide.shapes
        # Add title to page
        shapes.title.text = i
        
        #p = title.text_frame.paragraphs[0]
        #p.font.size = 64

        # Establish rows and columns of df; add +1 b/c it doesn't count index and column heads
        rows, cols = df.shape
        rows = rows+1
        cols = cols+1 

        # Size and placement of table
        left = Inches(.25)
        top = Inches(1)
        width = Inches(9.5)
        height = Inches(1.0)
    
        table = shapes.add_table(rows, cols, left, top, width, height).table
    
        # Write column headings into ppt table
        for i in range(cols-1):
            cell = table.cell(0, i+1)
            cell.text = list(df.columns.values)[i]
            paras = cell.text_frame.paragraphs
            for para in paras:
                para.font.bold = True
                para.font.size = Pt(12)
                para.font.name = 'Arial'
        
        # Create a list of lists of each column
        table_values = []
        for i in list(df.columns.values):
            table_values.append([n for n in df[i]])
    
        # Iterate over list of lists to populate ppt table
        for b in range(cols-1):
            for a in range(rows-1):
                cell = table.cell(a+1, b+1)
                cell.text = str(table_values[b][a])
                paras = cell.text_frame.paragraphs
                for para in paras:
                    para.font.bold = False
                    para.font.size = Pt(12)
                    para.font.name = 'Arial'

        # Add index values
        for a in range(rows-1):
            cell = table.cell(a+1, 0)
            cell.text = list(df.index.values)[a]
            #para = cell.text_frame.paragraphs[0]
            paras = cell.text_frame.paragraphs
            print(paras, len(paras))
            if len(paras) == 1:
                for para in paras:
                    para.font.bold = True
                    para.font.size = Pt(12)
                    para.font.name = 'Arial'
            if len(paras) == 3:
                    paras[0].font.bold = True
                    paras[0].font.size = Pt(12)
                    paras[0].font.name = 'Arial'
                    paras[1].font.bold = False
                    paras[1].font.size = Pt(12)
                    paras[1].font.name = 'Arial'
                    paras[2].font.bold = False
                    paras[2].font.size = Pt(12)
                    paras[2].font.name = 'Arial'
        
    prs.save(path)
    
    