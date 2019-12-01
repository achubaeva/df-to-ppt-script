from pptx import Presentation
from pptx.util import Inches

'''Create a powerpoint presentation with title page and following pages with tables from dataframe.'''

def create_title_page(title_str, subtitle_str):
    '''Title + subtitle page.'''
    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = title_str
    subtitle.text = subtitle_str

def create_new_page(title_str, df):
    '''Title and df to be converted into table.'''
    prs = Presentation()
    # Round numbers if necessary
    df = df.round(2)
    # First establish title only slide
    title_only_slide_layout = prs.slide_layouts[5]
    # Add slide to presentation
    slide = prs.slides.add_slide(title_only_slide_layout)
    shapes = slide.shapes
    # Add title to page
    shapes.title.text = title_str

    # Establish rows and columns of df; add +1 b/c it doesn't count index and column heads
    rows, cols = df.shape
    rows = rows+1
    cols = cols+1 

    # Size and placement of table
    left = Inches(.25)
    top = Inches(1.5)
    width = Inches(9.5)
    height = Inches(1.0)
    
    table = shapes.add_table(rows, cols, left, top, width, height).table
    
    # Write column headings into ppt table
    for i in range(cols-1):
        table.cell(0, i+1).text = list(df.columns.values)[i]
        
    # Create a list of lists of each column
    table_values = []
    for i in list(df.columns.values):
        table_values.append([n for n in df[i]])
    
    # Iterate over list of lists to populate ppt table
    for b in range(cols-1):
        for a in range(rows-1):
            table.cell(a+1, b+1).text = str(table_values[b][a])

    # Add index values
    for a in range(rows-1):
        table.cell(a+1, 0).text = list(df.index.values)[a]
